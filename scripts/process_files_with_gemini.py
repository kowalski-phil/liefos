#!/usr/bin/env python3
"""
Process any folder of files through Gemini (via OpenRouter) — extract content, summarise, analyse.

Usage:
  python scripts/process_files_with_gemini.py                    # processes demo_files/
  python scripts/process_files_with_gemini.py path/to/folder     # processes custom folder

Supported: PDF, PPTX, XLSX, DOCX, CSV, JSON, XML, MD, TXT, PY, JS, HTML, CSS, any text file

Output: outputs/file_summaries/YYYY-MM-DD/
  - <filename>_summary.md    (per-file analysis)
  - MASTER_SUMMARY.md        (full digest of all files)
"""

import os
import sys
import json
import csv
import io

# Fix Windows console encoding for emoji/unicode output
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

from pathlib import Path
from datetime import date
from dotenv import load_dotenv
from urllib.request import Request, urlopen
from urllib.error import HTTPError

load_dotenv()

API_KEY = os.environ.get("OPENROUTER_API_KEY")
if not API_KEY:
    print("Missing OPENROUTER_API_KEY in .env")
    sys.exit(1)

# ─────────────────────────────────────────────
# CONFIG — edit these to customise behaviour
# ─────────────────────────────────────────────
#
# MODEL: which model to use (via OpenRouter)
#   "google/gemini-3-flash-preview"  — fast, cheap, great for most files  ← default
#   "google/gemini-3-pro-preview"    — slower, higher quality on dense docs
#
# OUTPUT_DIR: where summaries are saved (default: outputs/file_summaries/YYYY-MM-DD/)
#
# max_chars (in analyse_with_gemini below): how much text to send per file
#   Default 12,000 chars. Raise for dense docs, lower to save API quota.
#
# ANALYSIS_PROMPT (below): the instruction the model follows for every file.
#   Tell Claude Code to edit it in plain English — e.g.:
#   "Add a Confidence score to every summary" or
#   "Always extract the author name and date if present"
# ─────────────────────────────────────────────

MODEL = "google/gemini-3-flash-preview"
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

BASE_DIR = Path(__file__).parent.parent
TODAY = date.today().isoformat()
OUTPUT_DIR = BASE_DIR / "outputs" / "file_summaries" / TODAY
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# ── Content Extractors ────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            text = "\n\n".join(p.extract_text() or "" for p in pdf.pages)
        return text.strip()
    except ImportError:
        pass
    try:
        from reportlab.lib.pagesizes import letter
        # fallback: read raw bytes, return notice
        return f"[PDF — {path.stat().st_size:,} bytes. Install pdfplumber for text extraction.]"
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name
                if "Heading" in style:
                    parts.append(f"\n## {para.text}")
                else:
                    parts.append(para.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    rows.append(" | ".join(str(c) if c is not None else "" for c in row))
            parts.extend(rows[:50])  # cap at 50 rows for context
            if ws.max_row > 50:
                parts.append(f"... ({ws.max_row - 50} more rows)")
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_csv(path: Path) -> str:
    try:
        with open(path, newline="", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        lines = [" | ".join(row) for row in rows[:30]]
        if len(rows) > 30:
            lines.append(f"... ({len(rows) - 30} more rows)")
        return "\n".join(lines)
    except Exception as e:
        return f"[CSV extraction error: {e}]"


def extract_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[Text extraction error: {e}]"


EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    ".csv":  extract_csv,
}

TEXT_EXTENSIONS = {".txt", ".md", ".json", ".xml", ".py", ".js", ".ts",
                   ".html", ".htm", ".css", ".yaml", ".yml", ".toml", ".sh"}


def extract_content(path: Path) -> tuple[str, str]:
    """Returns (content, file_type_label)."""
    ext = path.suffix.lower()

    if ext in EXTRACTORS:
        label = ext.lstrip(".").upper()
        return EXTRACTORS[ext](path), label
    elif ext in TEXT_EXTENSIONS:
        return extract_text(path), ext.lstrip(".").upper()
    else:
        # Try plain text anyway
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            return content, "TEXT"
        except Exception:
            return f"[Binary file — cannot extract text from {ext}]", "BINARY"


# ── OpenRouter Analyser ──────────────────────────────────────────────────────

ANALYSIS_PROMPT = """You are writing a quick-reference note for an Obsidian knowledge base.
The note will be scanned — not read. It needs to be instantly useful, conversational, and memorable.

File: {filename}
Type: {file_type}

Extracted content:
---
{content}
---

Silently classify the file as one of these types (do NOT output the classification — just use the correct format):

A) DELIVERABLE — invoice, report, contract, meeting notes, budget, presentation, sales data, proposal
B) REFERENCE — code, config, HTML, CSS, script, template, readme, design system, transcript, guide

Write the note using the matching format below. Output ONLY the note — no preamble, no classification label.

---

FORMAT A — DELIVERABLE FILES:

## TL;DR
One punchy sentence. What is this and what's the key number/date/outcome?

## Numbers & Dates
Only the figures that actually matter. No fluff.
- [key figure or deadline]
- [key figure or deadline]
- [key figure or deadline]

## What This Means
2-3 short, conversational sentences. Not bullet points — write like you're telling a colleague over Slack.
Focus on implications, not just facts. What's surprising, important, or worth flagging?

## Next Move
1-2 specific, non-obvious actions. Skip anything obvious like "open the file" or "read the document."
If nothing meaningful needs doing, omit this section entirely.

---

FORMAT B — REFERENCE FILES:

## TL;DR
One punchy sentence. What does this file do or contain?

## What's Inside
3-5 bullets. Be specific — name actual functions, config keys, sections, or design tokens. No vague descriptions.

## Worth Knowing
1-2 sentences on something non-obvious about how this file works or connects to other things.
If there's nothing worth flagging, omit this section.

---

Rules:
- Write conversationally. Short sentences. No consultant-speak.
- Never use the phrase "this file" more than once.
- No padding. If a section would be filler, leave it out.
- Total response under 250 words."""


def analyse_with_gemini(filename: str, file_type: str, content: str) -> str:
    # Truncate very long content
    max_chars = 12000
    if len(content) > max_chars:
        content = content[:max_chars] + f"\n\n[... truncated — original was {len(content):,} chars]"

    prompt = ANALYSIS_PROMPT.format(
        filename=filename,
        file_type=file_type,
        content=content
    )

    try:
        payload = json.dumps({
            "model": MODEL,
            "messages": [{"role": "user", "content": prompt}],
        }).encode("utf-8")

        req = Request(OPENROUTER_URL, data=payload, method="POST")
        req.add_header("Authorization", f"Bearer {API_KEY}")
        req.add_header("Content-Type", "application/json")

        with urlopen(req) as resp:
            data = json.loads(resp.read().decode("utf-8"))

        return data["choices"][0]["message"]["content"].strip()
    except HTTPError as e:
        body = e.read().decode("utf-8", errors="replace")
        return f"[OpenRouter error {e.code}: {body}]"
    except Exception as e:
        return f"[OpenRouter error: {e}]"


# ── Main Pipeline ─────────────────────────────────────────────────────────────

def process_folder(folder: Path):
    files = sorted([
        f for f in folder.iterdir()
        if f.is_file() and not f.name.startswith(".")
    ])

    if not files:
        print(f"No files found in {folder}")
        sys.exit(1)

    print(f"\n🔍 Processing {len(files)} files from: {folder}")
    print(f"📂 Output folder: {OUTPUT_DIR}\n")
    print("─" * 60)

    summaries = []

    for i, file_path in enumerate(files, 1):
        print(f"\n[{i}/{len(files)}] 📄 {file_path.name}")

        content, file_type = extract_content(file_path)
        char_count = len(content)
        print(f"  Type: {file_type} | Extracted: {char_count:,} chars")

        if "[Binary file" in content:
            print(f"  ⏭  Skipping binary file")
            continue

        print(f"  🤖 Sending to OpenRouter ({MODEL})...")
        analysis = analyse_with_gemini(file_path.name, file_type, content)

        # Save individual summary
        safe_name = file_path.stem.replace(" ", "_")
        out_path = OUTPUT_DIR / f"{safe_name}_summary.md"
        out_path.write_text(
            f"# {file_path.name}\n\n"
            f"**File type:** {file_type} | **Processed:** {TODAY}\n\n"
            f"{analysis}\n"
        )
        print(f"  ✅ Saved: {out_path.name}")

        summaries.append({
            "filename": file_path.name,
            "file_type": file_type,
            "chars": char_count,
            "analysis": analysis,
        })

    # Build master summary
    print("\n─" * 60)
    print("\n📋 Building master summary...")

    master_parts = [
        f"# File Processing Report\n",
        f"**Source folder:** `{folder}`  ",
        f"**Processed:** {TODAY}  ",
        f"**Files analysed:** {len(summaries)}\n",
        f"---\n",
    ]

    for s in summaries:
        master_parts.append(f"## {s['filename']} `{s['file_type']}`\n")
        master_parts.append(s["analysis"])
        master_parts.append("\n---\n")

    master_path = OUTPUT_DIR / "MASTER_SUMMARY.md"
    master_path.write_text("\n".join(master_parts))
    print(f"✅ Master summary: {master_path}")

    print(f"\n🎉 Done! Processed {len(summaries)} files.")
    print(f"   Output: {OUTPUT_DIR}")

    return OUTPUT_DIR


def main():
    if len(sys.argv) > 1:
        folder = Path(sys.argv[1])
    else:
        folder = BASE_DIR / "demo_files"

    if not folder.exists():
        print(f"Folder not found: {folder}")
        sys.exit(1)

    process_folder(folder)


if __name__ == "__main__":
    main()
