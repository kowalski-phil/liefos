#!/usr/bin/env python3
"""
Process Documents → Obsidian Markdown Notes
Uses Gemini Flash (via OpenRouter) to read any file type, extract signal, and output
clean compressed Markdown notes ready for Obsidian ingestion.

Usage:
    python process_docs_to_obsidian.py <input_folder> <output_folder>

Example:
    python process_docs_to_obsidian.py ~/Documents/company_files ~/vault/inbox

Supported file types: PDF, PPTX, DOCX, TXT, MD
"""

import os
import sys
import json

# Fix Windows console encoding for emoji/unicode output
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

from pathlib import Path
from datetime import date
from urllib.request import Request, urlopen
from urllib.error import HTTPError

from dotenv import load_dotenv
load_dotenv()

# ─────────────────────────────────────────────
# CONFIG — edit these to customise behaviour
# ─────────────────────────────────────────────
#
# MODEL: which model to use (via OpenRouter)
#   "google/gemini-3-flash-preview"  — fast, cheap, great for most files  ← default
#   "google/gemini-3-pro-preview"    — slower, higher quality for dense/complex docs
#
# SUPPORTED: file extensions to process — add or remove as needed
#
# SYNTHESIS_PROMPT (below): the instruction the model follows for EVERY file.
#   This is the most powerful thing to customise. Tell Claude Code:
#   "Edit the SYNTHESIS_PROMPT to focus on action items" or
#   "Add a Key People section that extracts names and roles"
#   and it will rewrite it for you.
# ─────────────────────────────────────────────

MODEL = "google/gemini-3-flash-preview"
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
API_KEY = os.environ.get("OPENROUTER_API_KEY")
TODAY = date.today().isoformat()

SUPPORTED = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".md"}

# ─────────────────────────────────────────────
# THE SYNTHESIS PROMPT
# Embedded in every API call — this is the brain of the script
# ─────────────────────────────────────────────

SYNTHESIS_PROMPT = """
You are processing a document to be stored in an Obsidian second brain vault.

Your job is to extract SIGNAL and discard NOISE.

SIGNAL = key insights, decisions, frameworks, facts, action items, names, dates that matter
NOISE  = headers/footers, legal boilerplate, filler sentences, repeated content, formatting artifacts

Output ONLY a clean Obsidian-compatible Markdown note with this exact structure:

---
type: [note | meeting | report | presentation | research | reference | other]
topic: [2-4 word topic description]
source: [FILENAME_PLACEHOLDER]
date_processed: [DATE_PLACEHOLDER]
tags: [comma-separated lowercase tags, no #]
---

# [Concise, descriptive title — what IS this document?]

## Key Insights
[3-7 bullet points of the highest-signal takeaways]

## Context
[1-2 sentences: what is this document, who created it, what was it for?]

## Details Worth Keeping
[Any specific data, frameworks, quotes, or facts that should be preserved verbatim]

## Action Items
[Only include if action items exist. Delete this section if none.]

---

RULES:
- Total note length: 300-600 words MAX. Ruthlessly compress.
- Do not copy paste large blocks of text. Synthesize.
- If the document is very short (under 200 words), keep it mostly intact.
- Write in clean, clear prose. No corporate speak.
- If you cannot read the document or it's empty, return: "# [filename] — Could not process"
"""

# ─────────────────────────────────────────────
# FILE READERS
# ─────────────────────────────────────────────

def read_pdf_text(file_path: Path) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            text = "\n\n".join(p.extract_text() or "" for p in pdf.pages)
        return text.strip()
    except ImportError:
        return f"[PDF — {file_path.stat().st_size:,} bytes. Install pdfplumber for text extraction.]"
    except Exception as e:
        return f"[PDF extraction error: {e}]"

def read_pptx_text(file_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes
                     if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                slides.append(f"[Slide {i}] " + " | ".join(texts))
        return "\n".join(slides)
    except ImportError:
        return "[Error: install python-pptx with: pip install python-pptx]"

def read_docx_text(file_path: Path) -> str:
    try:
        import docx
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return "[Error: install python-docx with: pip install python-docx]"

def read_text(file_path: Path) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

# ─────────────────────────────────────────────
# CORE PROCESSOR
# ─────────────────────────────────────────────

def process_file(file_path: Path) -> str | None:
    """Send a file to OpenRouter and return synthesized Markdown."""

    suffix = file_path.suffix.lower()
    prompt = SYNTHESIS_PROMPT.replace(
        "FILENAME_PLACEHOLDER", file_path.name
    ).replace(
        "DATE_PLACEHOLDER", TODAY
    )

    try:
        if suffix == ".pdf":
            text = read_pdf_text(file_path)
            content = f"PDF CONTENT:\n\n{text}\n\n{prompt}"

        elif suffix in {".pptx", ".ppt"}:
            text = read_pptx_text(file_path)
            content = f"PRESENTATION CONTENT:\n\n{text}\n\n{prompt}"

        elif suffix in {".docx", ".doc"}:
            text = read_docx_text(file_path)
            content = f"DOCUMENT CONTENT:\n\n{text}\n\n{prompt}"

        elif suffix in {".txt", ".md"}:
            text = read_text(file_path)
            content = f"TEXT CONTENT:\n\n{text}\n\n{prompt}"

        else:
            return None

        payload = json.dumps({
            "model": MODEL,
            "messages": [{"role": "user", "content": content}],
        }).encode("utf-8")

        req = Request(OPENROUTER_URL, data=payload, method="POST")
        req.add_header("Authorization", f"Bearer {API_KEY}")
        req.add_header("Content-Type", "application/json")

        with urlopen(req) as resp:
            data = json.loads(resp.read().decode("utf-8"))

        return data["choices"][0]["message"]["content"].strip()

    except HTTPError as e:
        body = e.read().decode("utf-8", errors="replace")
        print(f"  ✗ OpenRouter error {e.code}: {body}")
        return None
    except Exception as e:
        print(f"  ✗ Error: {e}")
        return None


# ─────────────────────────────────────────────
# BATCH RUNNER
# ─────────────────────────────────────────────

def process_folder(input_folder: str, output_folder: str):
    if not API_KEY:
        print("Error: OPENROUTER_API_KEY not set in .env")
        sys.exit(1)

    input_path = Path(input_folder).expanduser()
    output_path = Path(output_folder).expanduser()
    output_path.mkdir(parents=True, exist_ok=True)

    files = [f for f in input_path.iterdir() if f.suffix.lower() in SUPPORTED]

    if not files:
        print(f"No supported files found in {input_path}")
        print(f"Supported types: {', '.join(SUPPORTED)}")
        return

    print(f"\n📂 Input:  {input_path}")
    print(f"📁 Output: {output_path}")
    print(f"📄 Files:  {len(files)} found\n")
    print("─" * 50)

    success, skipped = 0, 0

    for i, file_path in enumerate(files, 1):
        print(f"[{i}/{len(files)}] {file_path.name}")

        result = process_file(file_path)

        if result:
            out_file = output_path / (file_path.stem + ".md")
            with open(out_file, "w", encoding="utf-8") as f:
                f.write(result)
            print(f"  ✓ → {out_file.name}")
            success += 1
        else:
            print(f"  ✗ Skipped")
            skipped += 1

    print("─" * 50)
    print(f"\n✅ Done: {success} notes created, {skipped} skipped")
    print(f"📁 Find your notes in: {output_path}")
    print("\nNext step: open Claude Code in your vault and run /vault-setup")
    print("to compartmentalize these notes into the right folders.")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    process_folder(sys.argv[1], sys.argv[2])
